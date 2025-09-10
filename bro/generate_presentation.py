#!/usr/bin/env python3
"""
Script para generar presentación PowerPoint sobre la Transformación Digital de Disney
Post-Adquisición de 21st Century Fox

Requiere: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def create_disney_presentation():
    """Crear presentación PowerPoint con 5 diapositivas sobre transformación digital Disney"""
    
    # Crear presentación
    prs = Presentation()
    
    # Colores Disney
    disney_blue = RGBColor(31, 85, 130)  # #1f5582
    disney_gold = RGBColor(255, 199, 44)  # #ffc72c
    white = RGBColor(255, 255, 255)
    dark_gray = RGBColor(64, 64, 64)
    
    # =================== SLIDE 1: TÍTULO ===================
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide1 = prs.slides.add_slide(slide_layout)
    
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Transformación Digital en Disney"
    title.text_frame.paragraphs[0].font.color.rgb = disney_blue
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Post-Adquisición de 21st Century Fox\n\nAnálisis del proceso de transformación digital más importante del sector entretenimiento"
    subtitle.text_frame.paragraphs[0].font.color.rgb = disney_gold
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[1].font.color.rgb = dark_gray
    subtitle.text_frame.paragraphs[1].font.size = Pt(18)
    
    # =================== SLIDE 2: CONTEXTO ADQUISICIÓN ===================
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide2 = prs.slides.add_slide(slide_layout)
    
    title = slide2.shapes.title
    title.text = "Contexto: La Gran Adquisición"
    title.text_frame.paragraphs[0].font.color.rgb = disney_blue
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    # Agregar contenido
    content = slide2.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Monto de la adquisición
    p = tf.paragraphs[0]
    p.text = "$71.3 mil millones"
    p.font.size = Pt(32)
    p.font.color.rgb = disney_gold
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Activos adquiridos
    p = tf.add_paragraph()
    p.text = "Activos Clave Adquiridos (2019):"
    p.font.size = Pt(20)
    p.font.color.rgb = disney_blue
    p.font.bold = True
    p.space_before = Pt(20)
    
    assets = [
        "🎭 20th Century Fox Studios",
        "📺 FX Networks", 
        "🌍 National Geographic",
        "🎮 Control total de Hulu"
    ]
    
    for asset in assets:
        p = tf.add_paragraph()
        p.text = asset
        p.font.size = Pt(18)
        p.space_before = Pt(12)
        p.level = 1
    
    # Objetivo
    p = tf.add_paragraph()
    p.text = "Objetivo: Consolidar contenido premium para competir en la era del streaming"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = dark_gray
    p.space_before = Pt(20)
    p.alignment = PP_ALIGN.CENTER
    
    # =================== SLIDE 3: ESTRATEGIA TRANSFORMACIÓN ===================
    slide_layout = prs.slide_layouts[1]
    slide3 = prs.slides.add_slide(slide_layout)
    
    title = slide3.shapes.title
    title.text = "Estrategia de Transformación Digital"
    title.text_frame.paragraphs[0].font.color.rgb = disney_blue
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide3.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Pilares fundamentales
    strategies = [
        {
            "title": "🎯 Direct-to-Consumer (DTC)",
            "items": ["Disney+", "ESPN+", "Hulu (control total)"]
        },
        {
            "title": "🔄 Integración de Plataformas", 
            "items": ["Unificación de sistemas", "Experiencia de usuario cohesiva", "Cross-platform synergy"]
        },
        {
            "title": "🤖 Data Analytics & AI",
            "items": ["Personalización de contenido", "Optimización de recomendaciones", "Machine Learning avanzado"]
        }
    ]
    
    for i, strategy in enumerate(strategies):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = strategy["title"]
        p.font.size = Pt(22)
        p.font.color.rgb = disney_blue
        p.font.bold = True
        if i > 0:
            p.space_before = Pt(20)
        
        for item in strategy["items"]:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.level = 1
            p.space_before = Pt(8)
    
    # =================== SLIDE 4: IMPLEMENTACIÓN TECNOLÓGICA ===================
    slide_layout = prs.slide_layouts[1]
    slide4 = prs.slides.add_slide(slide_layout)
    
    title = slide4.shapes.title
    title.text = "Implementación Tecnológica e IA"
    title.text_frame.paragraphs[0].font.color.rgb = disney_blue
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide4.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    tech_areas = [
        {
            "title": "☁️ AWS Partnership",
            "items": ["Escalabilidad global", "CDN optimizado", "Security Framework"]
        },
        {
            "title": "🧠 BAMTech Platform",
            "items": ["CRM Unificado (Vista 360° del cliente)", "Content Management centralizado", "Analytics con Machine Learning"]
        },
        {
            "title": "🔒 Protección Premium",
            "items": ["DRM avanzado", "Gestión compleja de licencias", "Frameworks de seguridad"]
        }
    ]
    
    for i, area in enumerate(tech_areas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = area["title"]
        p.font.size = Pt(20)
        p.font.color.rgb = disney_blue
        p.font.bold = True
        if i > 0:
            p.space_before = Pt(18)
        
        for item in area["items"]:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.level = 1
            p.space_before = Pt(6)
    
    # =================== SLIDE 5: RESULTADOS FINANCIEROS ===================
    slide_layout = prs.slide_layouts[1]
    slide5 = prs.slides.add_slide(slide_layout)
    
    title = slide5.shapes.title
    title.text = "Resultados Extraordinarios"
    title.text_frame.paragraphs[0].font.color.rgb = disney_blue
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide5.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Métricas principales
    p = tf.paragraphs[0]
    p.text = "📊 Métricas de Éxito"
    p.font.size = Pt(24)
    p.font.color.rgb = disney_blue
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    metrics = [
        "150M+ Suscriptores Disney+ (2023)",
        "Crecimiento más rápido que Netflix en sus primeros años",
        "45M+ Suscriptores Hulu (control total)",
        "Estrategia de bundling exitosa"
    ]
    
    for metric in metrics:
        p = tf.add_paragraph()
        p.text = metric
        p.font.size = Pt(18) if "150M+" in metric or "45M+" in metric else Pt(14)
        p.font.color.rgb = disney_gold if "150M+" in metric or "45M+" in metric else dark_gray
        p.font.bold = "150M+" in metric or "45M+" in metric
        p.space_before = Pt(12)
        p.level = 0 if "150M+" in metric or "45M+" in metric else 1
    
    # Logros clave
    p = tf.add_paragraph()
    p.text = "🏆 Logros Clave Alcanzados"
    p.font.size = Pt(20)
    p.font.color.rgb = disney_blue
    p.font.bold = True
    p.space_before = Pt(24)
    
    achievements = [
        "✅ Transformación B2B2C → B2C exitosa",
        "✅ Diversificación de ingresos consolidada", 
        "✅ Expansión internacional acelerada",
        "✅ Integración exitosa de activos Fox"
    ]
    
    for achievement in achievements:
        p = tf.add_paragraph()
        p.text = achievement
        p.font.size = Pt(16)
        p.space_before = Pt(8)
        p.level = 1
    
    # Lección clave
    p = tf.add_paragraph()
    p.text = "💡 \"La velocidad de implementación definió el éxito\""
    p.font.size = Pt(18)
    p.font.color.rgb = disney_gold
    p.font.bold = True
    p.font.italic = True
    p.space_before = Pt(20)
    p.alignment = PP_ALIGN.CENTER
    
    return prs

def main():
    """Función principal para generar la presentación"""
    try:
        print("🏰 Generando presentación PowerPoint: Transformación Digital Disney...")
        
        # Crear presentación
        prs = create_disney_presentation()
        
        # Guardar archivo
        output_path = "Disney_Transformacion_Digital_Presentacion.pptx"
        prs.save(output_path)
        
        print(f"✅ Presentación generada exitosamente: {output_path}")
        print(f"📁 Archivo guardado en: {os.path.abspath(output_path)}")
        print(f"📊 Total de diapositivas: {len(prs.slides)}")
        
        # Información adicional
        print("\n📋 Contenido de la presentación:")
        slide_titles = [
            "1. Título: Transformación Digital en Disney",
            "2. Contexto: La Gran Adquisición ($71.3B)",
            "3. Estrategia de Transformación Digital",
            "4. Implementación Tecnológica e IA", 
            "5. Resultados Extraordinarios y Métricas"
        ]
        
        for title in slide_titles:
            print(f"   {title}")
            
        print("\n🎯 Características:")
        print("   • Diseño profesional con colores Disney")
        print("   • Mínimo texto, máximo impacto visual")
        print("   • Enfoque en métricas y resultados clave")
        print("   • Iconos y elementos visuales integrados")
        
    except ImportError:
        print("❌ Error: La biblioteca 'python-pptx' no está instalada.")
        print("💡 Para instalarla, ejecuta: pip install python-pptx")
        return False
    except Exception as e:
        print(f"❌ Error al generar la presentación: {str(e)}")
        return False
    
    return True

if __name__ == "__main__":
    success = main()
    if success:
        print("\n🚀 ¡Presentación lista para usar!")
    else:
        print("\n💥 Error en la generación de la presentación.")